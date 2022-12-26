import pptx

file='a first format.pptx'

p=pptx.Presentation(file)
print(p.slide_width,p.slide_height,type(p.slide_width))

#for 4x3: w= 9144000 h=6858000
#for 16x9: w=9144000 h=5143500
#for 16x9 galyamin: w=12192000 h= 6858000
prs=p

bullet_slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
print(type(title_shape))
#body_shape = shapes.placeholders[1]

title_shape.text = 'Adding a Bullet Slide'