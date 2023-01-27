import pptx

file = 'sample.pptx'

p: pptx.Presentation = pptx.Presentation()
for (i, layout) in enumerate(p.slide_layouts):
    print(f'LAYOUT ---{i}-----')
    slide = p.slides.add_slide(layout)
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
