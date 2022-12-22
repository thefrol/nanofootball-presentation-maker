from pptx import Presentation
from pptx.util import Inches

p=Presentation()
title_only_slide_layput=p.slide_layouts[5]
slide= p.slides.add_slide(title_only_slide_layput)
shapes=slide.shapes

shapes.title.text='a python slide'

rows=cols=2

left=top=Inches(2.0)
width=Inches(6.0)
height=Inches(0.8)

table=shapes.add_table(rows,cols,left,top,width,height).table
table.cell(0,0).text='trener'
table.cell(0,1).text='takoyto'

table.cell(1,0).text='zadaniye'
table.cell(1,1).text='takoeto'

p.save('test.pptx')

