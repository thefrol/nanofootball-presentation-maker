
from pptx.slide import Slide
from pptx.util import Cm, Inches,Emu,Pt


from nf_presentation.builders.base import Builder,ElementBuilder,BoxElementBuilder,ColoredBoxBuilder
from nf_presentation.builders.paragraph_builder import ParagraphBuilder



from nf_presentation.logger import logger
from nf_presentation.settings import (DEFAULT_TABLE_WIDTH,
                       DEFAULT_TABLE_ROW_HEIGHT,
                       DEFAULT_TABLE_HORZ_BANDING,
                       DEFAULT_TABLE_MARK_FIRST_ROW,
                       DEFAULT_TEXT_FONT_SIZE_PT,
                       )




        

class RowDefinition:
    def __init__(self,texts,background=None,foreground=None,font_size_pt=DEFAULT_TEXT_FONT_SIZE_PT):
        self.background_color=background
        self.texts=texts
        self.foreground_color=foreground
        self.font_size_pt=font_size_pt

class RowTableBuilder(ElementBuilder):
    """a class building a table from rows,
    so we specify row height with row_height, and total width
    use append_row to add data with a tuples, or
    extend_rows with array of tuples with text
    
    the first rows wont be marked with background"""
    def __init__(self, width : float= DEFAULT_TABLE_WIDTH, row_height : float= DEFAULT_TABLE_ROW_HEIGHT, position : tuple[float,float]=(1,1)):
        self.rows=[]
        self.width=width
        self.row_height=row_height
        self.position=position



    def with_row_height(self, row_height):
        self.row_height=row_height
    
    def with_width(self, width):
        self.width=width
        return self

    def append_row(self, *data:list[str]):
        self.rows.append(tuple(data))
        return self

    def append_empty_row(self):
        return self.append_row()

    def extend_rows(self, rows_array:list[tuple]):
        self.rows.extend(rows_array)
        return self

    @property
    def cols_count(self):
        """returns amount of columns in a table, based on the data added by append_row and extend_rows
        so it returns a maximum amount of elements in a row"""
        if not self.rows:
            return 1  # the table is empty
        return max([len(row_items) for row_items in self.rows])

    @property
    def height(self):
        return self.row_height*self.rows_count

    @property
    def rows_count(self):
        return len(self.rows)

    

    def _build(self, slide:Slide):
        if self.rows_count==0:
            logger.info('the table is empty and not rendered')
            return
        table=slide.shapes.add_table(
            rows=self.rows_count,
            cols=self.cols_count,
            left=Cm(self.left),
            top=Cm(self.top),
            width=Cm(self.width),
            height=Cm(self.height)).table

        table.horz_banding=DEFAULT_TABLE_HORZ_BANDING
        table.first_row=DEFAULT_TABLE_MARK_FIRST_ROW  # disable coloring fist row


        for (row_number,row_content_list) in enumerate(self.rows):
            table.rows[row_number].height=Cm(self.row_height)
            initial_cell=table.cell(row_number,0)
            p=initial_cell.text_frame.paragraphs[0]
            p.font.size=Pt(DEFAULT_TEXT_FONT_SIZE_PT)  # if row is empty we need to set the font size anyways
                                                        #TODO: extract method

            last_filled_cell=initial_cell
            for (column_number, cell_content) in enumerate(row_content_list):
                #TODO if isinstance(cell_text,HTMLCell)
                #table.cell(row_number,column_number).text=str(cell_text)
                cell=table.cell(row_number,column_number)
                p=cell.text_frame.paragraphs[0]
                p.font.size=Pt(DEFAULT_TEXT_FONT_SIZE_PT)
                if isinstance(cell_content,str):
                    p.text=str(cell_content)
                elif isinstance(cell_content,int):
                    p.text=str(cell_content)
                elif isinstance(cell_content,float):
                    p.text=str(cell_content)
                elif cell_content is None:
                    logger.warning(f'Cell value is None. Replaced to ""')
                    p.text=''
                elif isinstance(cell_content,ParagraphBuilder):
                    cell_content:ParagraphBuilder=cell_content
                    cell_content._build(p)
                else:
                    logger.error(f'unknown content in a row of class {cell_content.__class__}.skiping cell')
                    continue
                
                
                last_filled_cell=cell
            #i this row is shorter than others we merge last cells
            try:              
                last_filled_column=len(row_content_list)
                if(last_filled_column<self.cols_count):
                    for empty_column_number in range(last_filled_column,self.cols_count):
                        empty_cell=table.cell(row_number,empty_column_number)
                        last_filled_cell.merge(empty_cell)
            except Exception:
                logger.error('Error while merging cells in a pptx table')