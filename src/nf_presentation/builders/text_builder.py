from pptx.slide import Slide
from pptx.util import Cm, Inches,Emu,Pt
from pptx.enum.text import PP_ALIGN

from nf_presentation.builders.base import Builder,ElementBuilder,BoxElementBuilder,ColoredBoxBuilder


from nf_presentation.logger import logger
from nf_presentation.settings import (
                       DEFAULT_TEXT_FONT_SIZE_PT,
                       DEFAULT_TITLE_FONT_SIZE_PT
                       )

class TextAlign:
    Center=PP_ALIGN.CENTER

class TextBuilder(ColoredBoxBuilder):
    def __init__(self, text:str,autosize=None, align=None, size_pt=DEFAULT_TEXT_FONT_SIZE_PT):
        super().__init__()
        self.text=text
        self.text_align=align
        self.autosize=autosize
        self.font_size_pt=size_pt
    def _build(self, slide: Slide):
        textbox=slide.shapes.add_textbox(
            left=Cm(self.left),
            top=Cm(self.top),
            width=Cm(self.width),
            height=Cm(self.height))

        text_frame=textbox.text_frame
        text_frame.auto_size=self.autosize
        text_frame.text=self.text

        p=text_frame.paragraphs[0]
        p.alignment=self.text_align

        font=p.font
            # font.name = 'Calibri'
            # font.size = Pt(18)
            # font.bold = True
        font.size=Pt(self.font_size_pt)
        if self.foreground_color:
            font.color.rgb = self.foreground_color

        if self.background_color:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb=self.background_color

        return textbox


class TitleBuilder(TextBuilder):
    def __init__(self, text: str):
        super().__init__(text, autosize=None,align=TextAlign.Center,size_pt=DEFAULT_TITLE_FONT_SIZE_PT)
    def _build(self, slide: Slide):
        textbox=super()._build(slide)
        slide.title=textbox