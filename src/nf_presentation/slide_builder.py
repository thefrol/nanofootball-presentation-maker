from abc import abstractmethod

from pptx import Presentation
from pptx.slide import Slide
from pptx.table import Table
from pptx.util import Cm

from .settings import DEFAULT_TABLE_WIDTH,DEFAULT_TABLE_ROW_HEIGHT

TITLE_ONLY_LAYOUT=5

class ElementBuilder:
    def __init__(self):
        pass
    @abstractmethod
    def _build(self, slide: Slide):
        pass

class Builder:
    @abstractmethod
    def _build(self, injection):
        pass

class RowTableBuilder(ElementBuilder):
    """a class building a table from rows,
    so we specify row height with row_height, and total width
    use append_row to add data with a tuples, or
    extend_rows with array of tuples with text
    
    the first rows wont be marked with background"""
    def __init__(self, width : float= DEFAULT_TABLE_WIDTH, row_height : float= DEFAULT_TABLE_ROW_HEIGHT, position : tuple[float,float]=(1,1)):
        self.rows=[]
        self.width=width
        self.row_height=row_height
        self.position=position
        

    @property
    def left(self):
        left,_ =self.position
        return left
    @property
    def top(self):
        _,top=self.position
        return top

 

    def append_row(self, *data:list[str]):
        self.rows.append(tuple(data))

    def extend_rows(self, rows_array:list[tuple]):
        self.rows.extend(rows_array)

    @property
    def cols_count(self):
        """returns amount of columns in a table, based on the data added by append_row and extend_rows
        so it returns a maximum amount of elements in a row"""
        return max(*[len(row_items) for row_items in self.rows])

    @property
    def height(self):
        return self.row_height*self.rows_count

    @property
    def rows_count(self):
        return len(self.rows)

    

    def _build(self, slide:Slide):
        table=slide.shapes.add_table(
            rows=self.rows_count,
            cols=self.cols_count,
            left=Cm(self.left),
            top=Cm(self.top),
            width=Cm(self.width),
            height=Cm(self.height)).table

        table.first_row=False  # disable coloring fist row

        for (row_number,row_text_list) in enumerate(self.rows):
            for (column_number, cell_text) in enumerate(row_text_list):
                #TODO if isinstance(cell_text,HTMLCell)
                table.cell(row_number,column_number).text=str(cell_text)


class SlideBuilder(Builder):
    """A builder for slides
    dont forget to add title
    and use any builder of table or picture to add to slides
    
    table_builder=TableBuilder()
    ...
    add_element(table_builder)"""
    def __init__(self,title=''):
        self.title=title
        self.shape_builders:list[ElementBuilder]=[]
    def set_title(self, title:str):
        self.title=title
        return self
    def add_element(self,builder:ElementBuilder):
        self.shape_builders.append(builder)
    def _build(self,presentation:Presentation):
        title_only_slide_layout=presentation.slide_layouts[TITLE_ONLY_LAYOUT]
        slide= presentation.slides.add_slide(title_only_slide_layout)

        shapes=slide.shapes
        shapes.title.text=self.title

        for builder in self.shape_builders:
            builder._build(slide)

        
    
class PresentationBuilder:
    def __init__(self):
        self.slide_builders:list[Builder]=[]

    def add_slide(self, slide_builder:SlideBuilder):
        self.slide_builders.append(slide_builder)

    def build(self) -> Presentation:
        presentation : Presentation =Presentation()
        for slide_builder in self.slide_builders:
            slide_builder._build(presentation)
        return presentation

    def save(self,to):
        """saves a presentation to a file or a stream
        to: file to save, 
        ex. save(to='sample.pptx')"""
        presentation=self.build()
        presentation.save(to)




