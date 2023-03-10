"""All sizes and lenghts are in centimeters (will be transformed to  pptx.util.Cm class)"""

from abc import abstractmethod
from dataclasses import dataclass
from typing import Union,Iterable,Iterator

from pptx import Presentation
from pptx.slide import Slide
from pptx.text.text import _Paragraph,_Run
from pptx.util import Cm, Inches,Emu,Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


from nf_presentation.logger import logger
from .settings import (DEFAULT_TABLE_WIDTH,
                       DEFAULT_TABLE_ROW_HEIGHT,
                       DEFAULT_TABLE_HORZ_BANDING,
                       DEFAULT_TABLE_MARK_FIRST_ROW,
                       DEFAULT_TEXT_FONT_SIZE_PT,
                       DEFAULT_TITLE_FONT_SIZE_PT
                       )


TITLE_ONLY_LAYOUT=5
BLANK_LAYOUT=6

class ElementBuilder:
    def __init__(self):
        self.position=(1,1)

    @property
    def left(self):
        left,_ =self.position
        return left

    @property
    def top(self):
        _,top=self.position
        return top

    def set_position(self,left,top):
        self.position=left,top
        return self

    def at(self, position):
        self.position=position
        return self

    @abstractmethod
    def _build(self, slide: Slide):
        pass

class BoxElementBuilder(ElementBuilder):
    def __init__(self):
        super().__init__()
        self.size=(1,1)
    
    def with_size(self, size:tuple[float,float]):
        self.size=size
        return self

    @property
    def width(self):
        width,height=self.size
        return width

    @property
    def height(self):
        width,height=self.size
        return height

class ColoredBoxBuilder(BoxElementBuilder):
    def __init__(self):
        super().__init__()
        self.background_color=None 
        self.foreground_color=None

    def with_foreground(self, rgb:tuple[int,int,int]):
        if rgb:
            r,g,b=rgb
            self.foreground_color=RGBColor(r,g,b)
        else:
            self.foreground_color=None
        return self

    def with_background(self, rgb:str):
        if rgb:
            r,g,b=rgb
            self.background_color=RGBColor(r,g,b)
        else:
            self.background_color=None
        return self


class ThemeColoredBoxBuilder(BoxElementBuilder):
    def __init__(self):
        super().__init__()
        self.background_color=None 
        self.foreground_color=None

    def with_foreground(self, color:str):
        if color:
            self.foreground_color=MSO_THEME_COLOR.from_xml(color)
        else:
            self.foreground_color=None
        return self

    def with_background(self, color:str):
        if color:
            self.background_color=MSO_THEME_COLOR.from_xml(color)
        else:
            self.background_color=None
        return self


class Builder:
    @abstractmethod
    def _build(self, injection):
        pass

class TextAlign:
    Center=PP_ALIGN.CENTER

class TextItem:
    def __init__(self):
        self.font_family=None
        self.font_size=None
    @abstractmethod
    def _make_run(self, run:_Run):
        pass

class PlainText(TextItem):
    def __init__(self,text:str):
        super().__init__()
        self.text=text
    def _make_run(self, run:_Run):
        super()._make_run(run)
        run.text=self.text

class HyperLink(TextItem):
    def __init__(self,link_text:str,href:str):
        super().__init__()
        self.link_text=link_text
        self.href=href
    def _make_run(self, run:_Run):
        super()._make_run(run)
        run.text=self.link_text
        run.hyperlink.address=self.href
        

class ParagraphBuilder(Builder):
    def __init__(self):
        self._items:list[TextItem]=[]
    def append_text(self,text:str,):
        self._items.append(PlainText(text=text))
    def append_link(self,link_text:str,href:str):
        self._items.append(HyperLink(link_text=link_text,href=href))
    def _build(self, paragraph:_Paragraph):
        for item in self._items:
            run=paragraph.add_run()
            item._make_run(run=run)

class TextBuilder(ColoredBoxBuilder):
    def __init__(self, text:str,autosize=None, align=None, size_pt=DEFAULT_TEXT_FONT_SIZE_PT):
        super().__init__()
        self.text=text
        self.text_align=align
        self.autosize=autosize
        self.font_size_pt=size_pt
    def _build(self, slide: Slide):
        textbox=slide.shapes.add_textbox(
            left=Cm(self.left),
            top=Cm(self.top),
            width=Cm(self.width),
            height=Cm(self.height))

        text_frame=textbox.text_frame
        text_frame.auto_size=self.autosize
        text_frame.text=self.text

        p=text_frame.paragraphs[0]
        p.alignment=self.text_align

        font=p.font
            # font.name = 'Calibri'
            # font.size = Pt(18)
            # font.bold = True
        font.size=Pt(self.font_size_pt)
        if self.foreground_color:
            font.color.rgb = self.foreground_color

        if self.background_color:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb=self.background_color

        return textbox


class TitleBuilder(TextBuilder):
    def __init__(self, text: str):
        super().__init__(text, autosize=None,align=TextAlign.Center,size_pt=DEFAULT_TITLE_FONT_SIZE_PT)
    def _build(self, slide: Slide):
        textbox=super()._build(slide)
        slide.title=textbox
        

class RowDefinition:
    def __init__(self,texts,background=None,foreground=None,font_size_pt=DEFAULT_TEXT_FONT_SIZE_PT):
        self.background_color=background
        self.texts=texts
        self.foreground_color=foreground
        self.font_size_pt=font_size_pt

class RowTableBuilder(ElementBuilder):
    """a class building a table from rows,
    so we specify row height with row_height, and total width
    use append_row to add data with a tuples, or
    extend_rows with array of tuples with text
    
    the first rows wont be marked with background"""
    def __init__(self, width : float= DEFAULT_TABLE_WIDTH, row_height : float= DEFAULT_TABLE_ROW_HEIGHT, position : tuple[float,float]=(1,1)):
        self.rows=[]
        self.width=width
        self.row_height=row_height
        self.position=position



    def with_row_height(self, row_height):
        self.row_height=row_height
    
    def with_width(self, width):
        self.width=width
        return self

    def append_row(self, *data:list[str]):
        self.rows.append(tuple(data))
        return self

    def append_empty_row(self):
        return self.append_row()

    def extend_rows(self, rows_array:list[tuple]):
        self.rows.extend(rows_array)
        return self

    @property
    def cols_count(self):
        """returns amount of columns in a table, based on the data added by append_row and extend_rows
        so it returns a maximum amount of elements in a row"""
        if not self.rows:
            return 1  # the table is empty
        return max([len(row_items) for row_items in self.rows])

    @property
    def height(self):
        return self.row_height*self.rows_count

    @property
    def rows_count(self):
        return len(self.rows)

    

    def _build(self, slide:Slide):
        if self.rows_count==0:
            logger.info('the table is empty and not rendered')
            return
        table=slide.shapes.add_table(
            rows=self.rows_count,
            cols=self.cols_count,
            left=Cm(self.left),
            top=Cm(self.top),
            width=Cm(self.width),
            height=Cm(self.height)).table

        table.horz_banding=DEFAULT_TABLE_HORZ_BANDING
        table.first_row=DEFAULT_TABLE_MARK_FIRST_ROW  # disable coloring fist row


        for (row_number,row_content_list) in enumerate(self.rows):
            table.rows[row_number].height=Cm(self.row_height)
            initial_cell=table.cell(row_number,0)
            p=initial_cell.text_frame.paragraphs[0]
            p.font.size=Pt(DEFAULT_TEXT_FONT_SIZE_PT)  # if row is empty we need to set the font size anyways
                                                        #TODO: extract method

            last_filled_cell=initial_cell
            for (column_number, cell_content) in enumerate(row_content_list):
                #TODO if isinstance(cell_text,HTMLCell)
                #table.cell(row_number,column_number).text=str(cell_text)
                cell=table.cell(row_number,column_number)
                p=cell.text_frame.paragraphs[0]
                p.font.size=Pt(DEFAULT_TEXT_FONT_SIZE_PT)
                if isinstance(cell_content,str):
                    p.text=str(cell_content)
                elif isinstance(cell_content,ParagraphBuilder):
                    cell_content:ParagraphBuilder=cell_content
                    cell_content._build(p)
                else:
                    logger.error(f'unknown content in a row of class {cell_content.__class__}.skiping cell')
                    continue
                
                
                last_filled_cell=cell
            #i this row is shorter than others we merge last cells
            try:              
                last_filled_column=len(row_content_list)
                if(last_filled_column<self.cols_count):
                    for empty_column_number in range(last_filled_column,self.cols_count):
                        empty_cell=table.cell(row_number,empty_column_number)
                        last_filled_cell.merge(empty_cell)
            except Exception:
                logger.error('Error while merging cells in a pptx table')


class ImageBuilder(BoxElementBuilder):
    """a class for adding images to a slide,
    automaticly scales for height or width
    so can use set_size(width,None) to scale image to width and same for height(leave width None)"""
    def __init__(self, image):
        super().__init__()
        self.image=image
        self.href=None

    def set_size(self,width=None,height=None):
        """DEPRECATED"""
        self.size=width,height
        return self

    def with_href(self, href:str):
        """adds a hyperlink to this shape, ex. builder.with_href('https://nanofootball.com')"""
        self.href=href
        return self

    @property
    def width(self):
        width,height=self.size
        return width
    @property
    def height(self):
        width,height=self.size
        return height

    def _build(self, slide: Slide):
        shape=slide.shapes.add_picture(
            image_file=self.image,
            left= Cm(self.left),
            top= Cm(self.top),
            width=Cm(self.width) if self.width else None,
            height=Cm(self.height) if self.height else None)
        if self.href:
            shape.click_action.hyperlink.address=self.href
    


class SlideBuilder(Builder):
    """A builder for slides
    dont forget to add title
    and use any builder of table or picture to add to slides
    
    table_builder=TableBuilder()
    ...
    add_element(table_builder)"""
    def __init__(self):
        self.shape_builders:list[ElementBuilder]=[]

    def set_title(self, title:str):
        self.title=title
        return self

    def add_element(self,builder:ElementBuilder):
        self.shape_builders.append(builder)
        return builder

    def create_table(self) -> RowTableBuilder:
        return self.add_element(RowTableBuilder())

    def create_image(self,image_file:str) -> ImageBuilder:
        return self.add_element(ImageBuilder(image=image_file))

    def create_text(self, text:str) -> TextBuilder:
        return self.add_element(TextBuilder(text=text))
    
    def create_title(self, text:str) -> TitleBuilder:
        return self.add_element(TitleBuilder(text=text))

    def _build(self,presentation:Presentation):
        blank_slide_layout=presentation.slide_layouts[BLANK_LAYOUT]
        slide= presentation.slides.add_slide(blank_slide_layout)

        for builder in self.shape_builders:
            builder._build(slide)

@dataclass
class PresentationRatio:
        width:Emu
        height:Emu
        alias:str

class Ratios:
    _16x9=PresentationRatio(width=Emu(12192000),height=Emu(6858000),alias='16:9')
    _4x3=PresentationRatio(width=Emu(9144000),height=Emu(6858000),alias='4:3') # the old broken one

    
    @classmethod
    def list(cls) -> Iterator[PresentationRatio]:
        """returns all ratios as iter"""
        ratios=[]
        for key in cls.__dict__:
            item:PresentationRatio=getattr(cls,key)
            if isinstance(item,PresentationRatio):
                ratios.append(item)
        return iter(ratios)


    @classmethod
    def get_by_alias(cls,alias:str):
        for ratio in cls.list():
            if ratio.alias==alias:
                return ratio
        raise AttributeError(f'Ratio {alias} not found')



    
class PresentationBuilder:
    def __init__(self, ratio : Union[str, PresentationRatio] = Ratios._16x9):
        self.slide_builders:list[Builder]=[]
        self.with_ratio(ratio)

    def with_ratio(self, ratio = Union[str, PresentationRatio]):
        if isinstance(ratio,str):
            self.ratio=Ratios.get_by_alias(ratio)
        else:
            self.ratio=ratio
        return self

    def add_slide(self, slide_builder:SlideBuilder):
        """DEPRECATED"""
        self.slide_builders.append(slide_builder)

    def create_slide(self) -> SlideBuilder:
        slide_builder=SlideBuilder()
        self.slide_builders.append(slide_builder)
        return slide_builder

    def build(self) -> Presentation:
        presentation : Presentation =Presentation()
        presentation.slide_width=self.ratio.width
        presentation.slide_height=self.ratio.height
        for slide_builder in self.slide_builders:
            slide_builder._build(presentation)
        return presentation

    def save(self,to):
        """saves a presentation to a file or a stream
        to: file to save, 
        ex. save(to='sample.pptx')"""
        presentation=self.build()
        presentation.save(to)




